# Generate a PPTX slide with syntax-highlighted JavaScript code using Pygments + python-pptx
# The file will be saved to /mnt/data/code_snippet.pptx
# You can download it after this cell runs.

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pygments import lex
from pygments.styles import get_style_by_name
from pygments.lexers import get_lexer_by_name

# Code snippet provided by the user
code = """var node = this.environment.variables.ros_node;
const publisher = node.createPublisher("geometry_msgs/msg/TwistStamped", "cmd_vel");
publisher.publish({
    header: {
        stamp: Date.now(),
        frame_id: "teleop_twist_joy"
    }, 
    twist: {
        linear: {x: speed, y: 0.0, z: 0.0},
        angular: {x: 0.0, y: 0.0, z: rot}
    }
});"""

# Settings
style_name = "vs"     # Pygments style ("default" possible)
font_name = "Cascadia code"     # Monospace font name (PowerPoint side)
font_size_pt = 17          # Font size
dark_bg = False             # Dark background for slide
line_numbers = True        # Show line numbers

# Prepare Pygments lexer and style
lexer = get_lexer_by_name("javascript")
style = get_style_by_name(style_name)

def hex_to_rgb(hexs):
    if not hexs:
        return (255, 255, 255)
    hexs = hexs.strip().lstrip('#')
    if len(hexs) == 6:
        return tuple(int(hexs[i:i+2], 16) for i in (0, 2, 4))
    if len(hexs) == 3:
        return tuple(int(c*2, 16) for c in hexs)
    return (255, 255, 255)

def run_style(run, style, tokentype):
    s = style.style_for_token(tokentype)
    if s.get('color'):
        r, g, b = hex_to_rgb(s['color'])
        run.font.color.rgb = RGBColor(r, g, b)
    if s.get('bold'):
        run.font.bold = True
    if s.get('italic'):
        run.font.italic = True

# Create PPTX
prs = Presentation()
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

# Background
bg = slide.background
fill = bg.fill
fill.solid()
if dark_bg:
    # Use style background if available; otherwise fallback to dark gray
    bg_hex = getattr(style, "background_color", None) or "1E2227"
    br, bgc, bb = hex_to_rgb(bg_hex)
    fill.fore_color.rgb = RGBColor(br, bgc, bb)
    default_fg = RGBColor(248, 248, 242)  # Monokai default foreground
else:
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    default_fg = RGBColor(33, 37, 41)

# Title
title_text = "Publish cmd_vel (JavaScript)"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1), Inches(0.6))
tf_title = title_box.text_frame
tf_title.text = title_text
p_title = tf_title.paragraphs[0]
p_title.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
p_title.font.size = Pt(font_size_pt + 6)
p_title.font.bold = True
p_title.font.color.rgb = default_fg

# Code textbox
left, top = Inches(0.5), Inches(1.0)
width, height = prs.slide_width - Inches(1), prs.slide_height - top - Inches(0.5)
box = slide.shapes.add_textbox(left, top, width, height)
tf = box.text_frame
tf.clear()
tf.word_wrap = True
tf.margin_left = Inches(0.25)
tf.margin_right = Inches(0.25)
tf.margin_top = Inches(0.15)
tf.margin_bottom = Inches(0.15)

# Initialize first paragraph
p = tf.paragraphs[0]
p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

# Line numbers helper
ln = 1
def add_ln_paragraph():
    global p, ln
    # create a new paragraph and (optionally) add the line number prefix
    para = tf.add_paragraph()
    para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    if line_numbers:
        rn = para.add_run()
        rn.text = f"{ln:>4} "
        rn.font.name = font_name
        rn.font.size = Pt(font_size_pt)
        rn.font.color.rgb = RGBColor(130, 130, 130)
    return para

# Add first line number if enabled
if line_numbers:
    rn0 = p.add_run()
    rn0.text = f"{ln:>4} "
    rn0.font.name = font_name
    rn0.font.size = Pt(font_size_pt)
    rn0.font.color.rgb = RGBColor(130, 130, 130)

# Lex and render tokens
tokens = list(lex(code, lexer))
for tokentype, val in tokens:
    parts = val.split('\n')
    for i, chunk in enumerate(parts):
        if i > 0:
            ln += 1
            p = add_ln_paragraph()
        if chunk == "":
            continue
        r = p.add_run()
        r.text = chunk
        r.font.name = font_name
        r.font.size = Pt(font_size_pt)
        r.font.color.rgb = default_fg
        run_style(r, style, tokentype)

# Save file
output_path = "code_snippet.pptx"
prs.save(output_path)
output_path
